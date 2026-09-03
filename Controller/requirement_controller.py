import os
import re
from werkzeug.utils import secure_filename
from flask import Blueprint, request, jsonify
from db import execute_query, execute_write
from agents.requirements_agent import normalize_requirements, analyze_conversational_intake
import json
from rag.vector_store import VectorStore
import PyPDF2
from io import BytesIO

requirement_bp = Blueprint('requirement', __name__)

def chunk_text(text, chunk_size=1500, overlap=200):
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

def derive_app_and_package_id(request_name, app_id=None, pkg_name=None):
    clean_slug = re.sub(r'[^a-zA-Z0-9]', '', request_name or '').lower() or 'app'
    default_id = f"com.company.{clean_slug}"
    
    final_app_id = app_id.strip() if (app_id and isinstance(app_id, str) and app_id.strip()) else default_id
    final_pkg_name = pkg_name.strip() if (pkg_name and isinstance(pkg_name, str) and pkg_name.strip()) else final_app_id
    
    return final_app_id, final_pkg_name

@requirement_bp.route('/', methods=['POST'])
def create_requirement():
    data = dict(request.form) if request.form else (request.json or {})
    
    # In the NLP flow, prompt is passed
    prompt = data.get('prompt', '')
    language = data.get('language', 'Java Kafka')

    file_uploads = request.files.getlist('file_upload')
    saved_file_paths = [] # This will now store [{"filename": "...", "content": "..."}]
    saved_image_paths = []
    
    for file_upload in file_uploads:
        if file_upload and file_upload.filename:
            filename = secure_filename(file_upload.filename)
            try:
                file_bytes = file_upload.read()
                content = ""
                filename_lower = filename.lower()
                
                if filename_lower.endswith(('.png', '.jpg', '.jpeg', '.webp')):
                    import base64
                    b64_content = base64.b64encode(file_bytes).decode('utf-8')
                    saved_image_paths.append(b64_content)
                elif filename_lower.endswith('.pdf'):
                    reader = PyPDF2.PdfReader(BytesIO(file_bytes))
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            content += page_text + "\n"
                elif filename_lower.endswith('.docx'):
                    import docx
                    doc = docx.Document(BytesIO(file_bytes))
                    content = "\n".join([para.text for para in doc.paragraphs])
                elif filename_lower.endswith('.pptx'):
                    from pptx import Presentation
                    prs = Presentation(BytesIO(file_bytes))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                else:
                    content = file_bytes.decode('utf-8', errors='ignore')
                saved_file_paths.append({"filename": filename, "content": content})
            except Exception as e:
                print(f"Error reading file {filename}: {e}")

    try:
        # Use LLM to extract ALL requirements from the unstructured prompt
        # We pass both the prompt and the file names to the LLM agent for context
        extract_payload = {
            'prompt': prompt,
            'language': language,
            'attached_files': [f"Filename: {f['filename']}\nContent:\n{f['content'][:5000]}" for f in saved_file_paths],
            'attached_images': saved_image_paths
        }
        normalized_spec = normalize_requirements(extract_payload)
        
        req_name = normalized_spec.get('request_name') or 'NLP Chat Request'
        final_app_id, final_pkg_name = derive_app_and_package_id(
            req_name, 
            normalized_spec.get('application_id'), 
            normalized_spec.get('package_name')
        )
        normalized_spec['application_id'] = final_app_id
        normalized_spec['package_name'] = final_pkg_name

        # Save request with track mapping
        ensure_request_track_columns()
        request_id = execute_write(
            "INSERT INTO generation_requests (request_name, application_id, package_name, requested_by, track_id, track_name, project_id) VALUES (%s, %s, %s, %s, %s, %s, %s)",
            (req_name, 
             final_app_id, 
             final_pkg_name, 
             data.get('requested_by', 'User'),
             data.get('track_id'),
             data.get('track_name'),
             data.get('project_id'))
        )
        
        if not request_id:
            return jsonify({"success": False, "message": "Failed to save request"}), 500
            
        # Save spec
        # Assuming sample_file_path is a text/varchar column, we store JSON array of filenames.
        paths_to_store = [f['filename'] for f in saved_file_paths]
        paths_str = json.dumps(paths_to_store) if paths_to_store else None
        
        execute_write(
            """INSERT INTO generation_specs (request_id, source_topics, target_topics, consumer_group, state_store_needed, error_topic_policy, schema_hints, sample_file_path, normalized_by) 
               VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)""",
            (request_id, normalized_spec.get('source_topics'), normalized_spec.get('target_topics'), 
             normalized_spec.get('consumer_group'), normalized_spec.get('state_store_needed', False), 
             normalized_spec.get('error_topic_policy', 'DLQ'), prompt, paths_str, 'ai')
        )
        
        # Add to Vector Store
        try:
            vs = VectorStore()
            docs = []
            if prompt.strip():
                docs.append(prompt)
            for f in saved_file_paths:
                if f['content'].strip():
                    chunks = chunk_text(f['content'])
                    for idx, chunk in enumerate(chunks):
                        docs.append(f"File {f['filename']} (Part {idx+1}):\n{chunk}")
            
            if docs:
                track_id_val = -1
                try:
                    if data.get('track_id'):
                        track_id_val = int(data.get('track_id'))
                except ValueError:
                    pass
                metas = [{"request_id": request_id, "type": "intake", "track_id": track_id_val} for _ in docs]
                ids = [f"req_{request_id}_intake_{i}" for i in range(len(docs))]
                vs.add_documents(docs, metas, ids)
        except Exception as e:
            print(f"Failed to add to VectorStore: {e}")
        
        return jsonify({"success": True, "data": {"request_id": request_id, "spec": normalized_spec}}), 201
    except ValueError as e:
        return jsonify({"success": False, "message": str(e)}), 400
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "message": str(e)}), 500

@requirement_bp.route('/intake-chat', methods=['POST'])
def intake_chat():
    data = dict(request.form) if request.form else (request.json or {})
    messages_str = data.get('messages', '[]')
    try:
        messages = json.loads(messages_str)
    except Exception:
        messages = []
        
    language = data.get('language', 'Java Kafka')
    
    if not messages:
        return jsonify({"success": False, "message": "No messages provided"}), 400
        
    file_uploads = request.files.getlist('file_upload')
    saved_file_paths = []
    saved_image_paths = []
    
    for file_upload in file_uploads:
        if file_upload and file_upload.filename:
            filename = secure_filename(file_upload.filename)
            try:
                file_bytes = file_upload.read()
                content = ""
                filename_lower = filename.lower()
                
                if filename_lower.endswith(('.png', '.jpg', '.jpeg', '.webp')):
                    import base64
                    b64_content = base64.b64encode(file_bytes).decode('utf-8')
                    saved_image_paths.append(b64_content)
                elif filename_lower.endswith('.pdf'):
                    reader = PyPDF2.PdfReader(BytesIO(file_bytes))
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            content += page_text + "\n"
                elif filename_lower.endswith('.docx'):
                    import docx
                    doc = docx.Document(BytesIO(file_bytes))
                    content = "\n".join([para.text for para in doc.paragraphs])
                elif filename_lower.endswith('.pptx'):
                    from pptx import Presentation
                    prs = Presentation(BytesIO(file_bytes))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                else:
                    content = file_bytes.decode('utf-8', errors='ignore')
                saved_file_paths.append({"filename": filename, "content": content})
            except Exception as e:
                print(f"Error reading file {filename}: {e}")
            
    try:
        files_content = [f"Filename: {f['filename']}\nContent:\n{f['content'][:5000]}" for f in saved_file_paths]
        result = analyze_conversational_intake(messages, language, files_content, images=saved_image_paths)
        
        if result.get('status') == 'complete':
            req = result.get('requirements', {})
            req_name = req.get('request_name') or 'NLP Chat Request'
            final_app_id, final_pkg_name = derive_app_and_package_id(
                req_name, 
                req.get('application_id'), 
                req.get('package_name')
            )
            req['application_id'] = final_app_id
            req['package_name'] = final_pkg_name

            # Save request with track mapping
            ensure_request_track_columns()
            request_id = execute_write(
                "INSERT INTO generation_requests (request_name, application_id, package_name, requested_by, track_id, track_name, project_id) VALUES (%s, %s, %s, %s, %s, %s, %s)",
                (req_name, 
                 final_app_id, 
                 final_pkg_name, 
                 'User',
                 data.get('track_id'),
                 data.get('track_name'),
                 data.get('project_id'))
            )
            
            # Save spec
            paths_to_store = [f['filename'] for f in saved_file_paths]
            paths_str = json.dumps(paths_to_store) if paths_to_store else None
            execute_write(
                """INSERT INTO generation_specs (request_id, source_topics, target_topics, consumer_group, state_store_needed, error_topic_policy, schema_hints, sample_file_path, normalized_by) 
                   VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)""",
                (request_id, req.get('source_topics'), req.get('target_topics'), 
                 req.get('consumer_group'), req.get('state_store_needed', False), 
                 req.get('error_topic_policy', 'DLQ'), json.dumps(messages), paths_str, 'ai')
            )
            
            # Add to Vector Store
            try:
                vs = VectorStore()
                docs = []
                chat_text = "\n".join([f"{msg.get('role')}: {msg.get('text')}" for msg in messages])
                if chat_text.strip():
                    docs.append(chat_text)
                for f in saved_file_paths:
                    if f['content'].strip():
                        chunks = chunk_text(f['content'])
                        for idx, chunk in enumerate(chunks):
                            docs.append(f"File {f['filename']} (Part {idx+1}):\n{chunk}")
                
                if docs:
                    track_id_val = -1
                    try:
                        if data.get('track_id'):
                            track_id_val = int(data.get('track_id'))
                    except ValueError:
                        pass
                    metas = [{"request_id": request_id, "type": "intake", "track_id": track_id_val} for _ in docs]
                    ids = [f"req_{request_id}_intake_{i}" for i in range(len(docs))]
                    vs.add_documents(docs, metas, ids)
            except Exception as e:
                print(f"Failed to add to VectorStore: {e}")
            
            return jsonify({"success": True, "status": "complete", "data": {"request_id": request_id}})
        else:
            return jsonify({"success": True, "status": "more_info", "question": result.get('question', "Can you provide more details?")})
            
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "message": str(e)}), 500

def ensure_request_track_columns():
    try:
        cols = execute_query("SHOW COLUMNS FROM generation_requests")
        col_names = [c['Field'] for c in cols] if cols else []
        
        if 'track_id' not in col_names:
            execute_write("ALTER TABLE generation_requests ADD COLUMN track_id INT")
        if 'track_name' not in col_names:
            execute_write("ALTER TABLE generation_requests ADD COLUMN track_name VARCHAR(255)")
        if 'project_id' not in col_names:
            execute_write("ALTER TABLE generation_requests ADD COLUMN project_id INT")
    except Exception as e:
        pass

@requirement_bp.route('/', methods=['GET'])
def list_requirements():
    ensure_request_track_columns()
    status = request.args.get('status')
    track_id = request.args.get('track_id')
    query = """
        SELECT r.*, s.schema_hints,
               (SELECT is_auto_approved FROM blueprints WHERE request_id = r.id ORDER BY id DESC LIMIT 1) as is_auto_approved,
               (SELECT job_status FROM pipeline_jobs WHERE request_id = r.id ORDER BY id DESC LIMIT 1) as job_status,
               (SELECT current_step FROM pipeline_jobs WHERE request_id = r.id ORDER BY id DESC LIMIT 1) as current_step,
               (SELECT COUNT(*) FROM validation_results WHERE request_id = r.id AND severity = 'error' AND passed = 0) as error_count
        FROM generation_requests r 
        LEFT JOIN generation_specs s ON r.id = s.request_id
        WHERE 1=1
    """
    params = []
    if status:
        query += " AND r.status=%s"
        params.append(status)
    if track_id:
        query += " AND r.track_id=%s"
        params.append(track_id)
        
    query += " ORDER BY r.created_at DESC"
    reqs = execute_query(query, tuple(params))
    return jsonify({"success": True, "data": reqs})

@requirement_bp.route('/<int:req_id>', methods=['GET'])
def get_requirement(req_id):
    req = execute_query("SELECT * FROM generation_requests WHERE id=%s", (req_id,))
    if not req:
        return jsonify({"success": False, "message": "Not found"}), 404
    spec = execute_query("SELECT * FROM generation_specs WHERE request_id=%s", (req_id,))
    bp = execute_query("SELECT * FROM blueprints WHERE request_id=%s ORDER BY id DESC LIMIT 1", (req_id,))
    job = execute_query("SELECT * FROM pipeline_jobs WHERE request_id=%s ORDER BY id DESC LIMIT 1", (req_id,))
    val_rows = execute_query("SELECT * FROM validation_results WHERE request_id=%s", (req_id,))
    
    has_errors = any(v['severity'] == 'error' and not v['passed'] for v in val_rows) if val_rows else False
    has_warnings = any(v['severity'] == 'warning' and not v['passed'] for v in val_rows) if val_rows else False

    if bp:
        req[0]['is_auto_approved'] = bp[0]['is_auto_approved']

    return jsonify({
        "success": True, 
        "data": {
            "request": req[0], 
            "spec": spec[0] if spec else None,
            "blueprint": bp[0] if bp else None,
            "job": job[0] if job else None,
            "validation_summary": {
                "results": val_rows,
                "has_errors": has_errors,
                "has_warnings": has_warnings,
                "count": len(val_rows)
            }
        }
    })

