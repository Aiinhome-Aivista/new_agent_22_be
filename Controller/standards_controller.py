import os
import subprocess
from flask import Blueprint, request, jsonify
from config import GEMINI_API_KEY
import google.generativeai as genai
from db import execute_query, execute_write

if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

standards_bp = Blueprint('standards', __name__)

KB_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'knowledge_base')

def trigger_ingest():
    ingest_script = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'rag', 'ingest.py')
    try:
        subprocess.run(['python', ingest_script], check=True)
        return True
    except Exception as e:
        print(f"Error running ingest: {e}")
        return False

def ensure_standards_track_columns():
    try:
        cols = execute_query("SHOW COLUMNS FROM architecture_standards")
        col_names = [c['Field'] for c in cols] if cols else []
        if 'track_id' not in col_names:
            execute_write("ALTER TABLE architecture_standards ADD COLUMN track_id INT")
        if 'track_name' in col_names:
            execute_write("ALTER TABLE architecture_standards DROP COLUMN track_name")
        if 'folder' not in col_names:
            execute_write("ALTER TABLE architecture_standards ADD COLUMN folder VARCHAR(255) DEFAULT 'standards'")
        if 'created_by' not in col_names:
            execute_write("ALTER TABLE architecture_standards ADD COLUMN created_by VARCHAR(255)")

        # Ensure created_by is never NULL in database
        execute_write("UPDATE architecture_standards SET created_by = '1' WHERE created_by IS NULL OR created_by = ''")
    except Exception as e:
        print(f"Error ensuring standards columns: {e}")

@standards_bp.route('/', methods=['GET'])
def list_standards():
    ensure_standards_track_columns()
    raw_track_id = request.args.get('track_id')
    track_id = int(raw_track_id) if raw_track_id and str(raw_track_id).isdigit() else None
    
    # 1. Fetch DB standards for this track (including fallback global standards)
    if track_id:
        db_standards = execute_query("SELECT * FROM architecture_standards WHERE track_id = %s OR track_id IS NULL ORDER BY id DESC", (track_id,))
    else:
        db_standards = execute_query("SELECT * FROM architecture_standards ORDER BY id DESC")
    
    # 2. If no DB standards exist at all, populate default standards into DB
    if not db_standards and os.path.exists(KB_DIR):
        for root, dirs, files in os.walk(KB_DIR):
            for file in files:
                if file.endswith('.md'):
                    file_path = os.path.join(root, file)
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                    title = file.replace('.md', '').replace('_', ' ').replace('-', ' ').title()
                    folder_name = os.path.basename(root)
                    if folder_name not in ['standards', 'miro_diagram', 'validation_rules', 'sample_scripts']:
                        folder_name = 'standards'
                    
                    execute_write(
                        "INSERT INTO architecture_standards (title, description, folder, created_by, track_id) VALUES (%s, %s, %s, %s, %s)",
                        (title, content, folder_name, '1', track_id)
                    )
        if track_id:
            db_standards = execute_query("SELECT * FROM architecture_standards WHERE track_id = %s OR track_id IS NULL ORDER BY id DESC", (track_id,))
        else:
            db_standards = execute_query("SELECT * FROM architecture_standards ORDER BY id DESC")

    standards = []
    if db_standards:
        for item in db_standards:
            raw_title = item.get('title', 'Standard') or 'Standard'
            filename_val = raw_title if raw_title.lower().endswith('.md') else f"{raw_title}.md"
            folder_val = item.get('folder', 'standards') or 'standards'
            if folder_val not in ['standards', 'miro_diagram', 'validation_rules', 'sample_scripts']:
                folder_val = 'standards'

            standards.append({
                "id": str(item['id']),
                "db_id": item['id'],
                "filename": filename_val,
                "folder": folder_val,
                "content": item.get('description', '') or '',
                "track_id": item.get('track_id'),
                "created_by": item.get('created_by') or '1'
            })

    return jsonify({"success": True, "data": standards})

@standards_bp.route('/', methods=['POST'])
def save_standard():
    ensure_standards_track_columns()
    data = request.json or {}
    filename = data.get('filename', '').strip()
    folder = data.get('folder', 'standards')
    content = data.get('content', '')
    raw_track_id = data.get('track_id')
    track_id = int(raw_track_id) if raw_track_id and str(raw_track_id).isdigit() else None
    created_by = str(data.get('created_by') or '1')
    
    if not filename:
        return jsonify({"success": False, "message": "Filename is required"}), 400
        
    is_edit = data.get('is_edit', False)
    raw_name = filename[:-3] if filename.lower().endswith('.md') else filename
    clean_title = raw_name.replace('_', ' ').replace('-', ' ').title()
    clean_filename = f"{raw_name}.md"

    try:
        if track_id:
            existing = execute_query("SELECT id FROM architecture_standards WHERE (title=%s OR title=%s) AND track_id=%s", (clean_filename, clean_title, track_id))
        else:
            existing = execute_query("SELECT id FROM architecture_standards WHERE (title=%s OR title=%s) AND track_id IS NULL", (clean_filename, clean_title))
            
        if existing and not is_edit:
            return jsonify({
                "success": False, 
                "message": f"A standard with filename '{clean_filename}' already exists. Please choose a different filename."
            }), 400

        if not track_id:
            target_dir = os.path.join(KB_DIR, folder)
            os.makedirs(target_dir, exist_ok=True)
            
            file_path = os.path.join(target_dir, clean_filename)
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content)

        if existing:
            execute_write(
                "UPDATE architecture_standards SET title=%s, description=%s, folder=%s, created_by=%s, track_id=%s WHERE id=%s", 
                (clean_title, content, folder, created_by, track_id, existing[0]['id'])
            )
        else:
            execute_write(
                "INSERT INTO architecture_standards (title, description, folder, created_by, track_id) VALUES (%s, %s, %s, %s, %s)", 
                (clean_title, content, folder, created_by, track_id)
            )
    except Exception as db_err:
        print(f"DB save sync error: {db_err}")
        return jsonify({"success": False, "message": str(db_err)}), 500
        
    trigger_ingest()
    
    return jsonify({"success": True, "message": "Standard saved successfully"})

@standards_bp.route('/upload', methods=['POST'])
def upload_standard():
    ensure_standards_track_columns()
    if 'file' not in request.files:
        return jsonify({"success": False, "message": "No file uploaded"}), 400
        
    file = request.files['file']
    if not file or file.filename == '':
        return jsonify({"success": False, "message": "No selected file"}), 400
        
    folder = request.form.get('folder', 'standards')
    raw_track_id = request.form.get('track_id')
    track_id = int(raw_track_id) if raw_track_id and str(raw_track_id).isdigit() else None
    created_by = str(request.form.get('created_by') or '1')
    
    filename = file.filename.strip()
    raw_name = filename[:-3] if filename.lower().endswith('.md') else filename
    clean_title = raw_name.replace('_', ' ').replace('-', ' ').title()
    clean_filename = f"{raw_name}.md"

    try:
        content = file.read().decode('utf-8', errors='ignore')
    except Exception as e:
        content = f"Uploaded content for {filename}"

    if not track_id:
        target_dir = os.path.join(KB_DIR, folder)
        os.makedirs(target_dir, exist_ok=True)
        
        file_path = os.path.join(target_dir, clean_filename)
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content)

    new_id = None
    try:
        if track_id:
            existing = execute_query("SELECT id FROM architecture_standards WHERE (title=%s OR title=%s) AND track_id=%s", (clean_filename, clean_title, track_id))
        else:
            existing = execute_query("SELECT id FROM architecture_standards WHERE (title=%s OR title=%s) AND track_id IS NULL", (clean_filename, clean_title))
            
        if existing:
            new_id = existing[0]['id']
            execute_write(
                "UPDATE architecture_standards SET title=%s, description=%s, folder=%s, created_by=%s, track_id=%s WHERE id=%s", 
                (clean_title, content, folder, created_by, track_id, new_id)
            )
        else:
            new_id = execute_write(
                "INSERT INTO architecture_standards (title, description, folder, created_by, track_id) VALUES (%s, %s, %s, %s, %s)", 
                (clean_title, content, folder, created_by, track_id)
            )
    except Exception as db_err:
        print(f"DB upload sync error: {db_err}")
        
    trigger_ingest()
    
    return jsonify({
        "success": True, 
        "message": "File uploaded successfully",
        "data": {
            "id": str(new_id) if new_id else clean_filename,
            "filename": clean_filename,
            "folder": folder,
            "content": content,
            "track_id": track_id,
            "created_by": created_by
        }
    })

def extract_text_from_stream(file_stream, filename):
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == '.pdf':
        err_msg = ""
        # 1. Try pypdf
        try:
            import pypdf
            reader = pypdf.PdfReader(file_stream)
            pages = []
            for i, page in enumerate(reader.pages):
                txt = page.extract_text() or ''
                if txt.strip():
                    pages.append(f"--- Page {i+1} ---\n{txt.strip()}")
            if pages:
                return "\n\n".join(pages)
        except Exception as e:
            err_msg += f"pypdf: {e}; "

        # 2. Try PyPDF2
        try:
            file_stream.seek(0)
            import PyPDF2
            reader = PyPDF2.PdfReader(file_stream)
            pages = []
            for i, page in enumerate(reader.pages):
                txt = page.extract_text() or ''
                if txt.strip():
                    pages.append(f"--- Page {i+1} ---\n{txt.strip()}")
            if pages:
                return "\n\n".join(pages)
        except Exception as e:
            err_msg += f"PyPDF2: {e}; "

        # 3. Try pdfplumber
        try:
            file_stream.seek(0)
            import pdfplumber
            with pdfplumber.open(file_stream) as pdf:
                pages = []
                for i, page in enumerate(pdf.pages):
                    txt = page.extract_text() or ''
                    if txt.strip():
                        pages.append(f"--- Page {i+1} ---\n{txt.strip()}")
                if pages:
                    return "\n\n".join(pages)
        except Exception as e:
            err_msg += f"pdfplumber: {e}; "

        return f"PDF Extraction Notice: No extractable text found in PDF ({err_msg})."

    elif ext in ['.docx', '.doc']:
        try:
            import docx
            doc = docx.Document(file_stream)
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            
            # Extract tables text if present
            table_lines = []
            for table in doc.tables:
                for row in table.rows:
                    row_txt = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_txt:
                        table_lines.append(row_txt)
            if table_lines:
                paragraphs.append("\n### Document Tables:\n" + "\n".join(table_lines))

            return "\n\n".join(paragraphs) if paragraphs else "No extractable text found in Word document."
        except Exception as e:
            return f"Word (.docx) Extraction Error: {str(e)}"

    elif ext in ['.pptx', '.ppt']:
        try:
            import pptx
            prs = pptx.Presentation(file_stream)
            slides = []
            for i, slide in enumerate(prs.slides):
                stext = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        stext.append(shape.text.strip())
                if stext:
                    slides.append(f"--- Slide {i+1} ---\n" + "\n".join(stext))
            return "\n\n".join(slides) if slides else "No extractable text found in PowerPoint presentation."
        except Exception as e:
            return f"PowerPoint Extraction Error: {str(e)}"

    elif ext in ['.xlsx', '.xls', '.csv']:
        try:
            import pandas as pd
            if ext == '.csv':
                df = pd.read_csv(file_stream)
                return df.to_markdown(index=False) if hasattr(df, 'to_markdown') else df.to_string(index=False)
            else:
                excel_file = pd.ExcelFile(file_stream)
                sheets = []
                for sname in excel_file.sheet_names:
                    df = pd.read_excel(excel_file, sheet_name=sname)
                    tbl = df.to_markdown(index=False) if hasattr(df, 'to_markdown') else df.to_string(index=False)
                    sheets.append(f"### Sheet: {sname}\n\n{tbl}")
                return "\n\n".join(sheets)
        except Exception as e:
            return f"Excel Extraction Error: {str(e)}"

    elif ext in ['.png', '.jpg', '.jpeg', '.webp']:
        if not GEMINI_API_KEY:
            return "Image Extraction Error: GEMINI_API_KEY is not configured in .env file."
        try:
            import google.generativeai as genai
            from PIL import Image
            
            genai.configure(api_key=GEMINI_API_KEY)
            img = Image.open(file_stream)
            model = genai.GenerativeModel('gemini-3.5-flash')
            prompt = """
            You are a Principal Enterprise Architect at PwC. Analyze this architecture or workflow diagram (like a Miro board).
            Your goal is to extract the content and convert it into a highly professional, production-ready "Architecture Standard & Blueprint" document.
            
            Please strictly follow this structure:
            
            # Enterprise Architecture Blueprint
            
            ## 1. Architectural Overview & Component Naming Conventions
            - Describe the overall architecture pattern (e.g., Event-Driven, Microservices, API-led).
            - List all components clearly with the following naming convention: `[Component Name] ([Technology/Role])`.
            
            ## 2. Standard Architectural Rules
            - Deduce strict architectural rules based on the diagram (e.g., "The UI must never connect directly to databases", "All inter-service communication must be asynchronous via Kafka").
            
            ## 3. Step-by-Step Process Flow
            - Create a beautiful ASCII diagram of the step-by-step flow.
            - Detail the exact Step-by-Step process flow, explicitly mentioning the Triggers, Actions, and Protocols (e.g., HTTP/REST, Async Publish, Consume).
            
            Use clean, professional Markdown. The output must be 100% ready for an enterprise engineering team to follow.
            """
            response = model.generate_content([prompt, img])
            return response.text
        except Exception as e:
            return f"Image Extraction Error: {str(e)}"

    else:
        try:
            return file_stream.read().decode('utf-8', errors='ignore')
        except Exception as e:
            return f"Error reading text file: {str(e)}"

@standards_bp.route('/parse-file', methods=['POST'])
def parse_uploaded_file():
    files = request.files.getlist('file') or request.files.getlist('files')
    if not files or all(f.filename == '' for f in files):
        return jsonify({"success": False, "message": "No files uploaded"}), 400

    extracted_items = []
    combined_texts = []
    
    for file in files:
        if file and file.filename:
            filename = file.filename.strip()
            raw_name = os.path.splitext(filename)[0]
            clean_filename = f"{raw_name}.md"
            extracted_text = extract_text_from_stream(file.stream, filename)
            
            extracted_items.append({
                "filename": clean_filename,
                "original_filename": filename,
                "content": extracted_text
            })
            combined_texts.append(f"<!-- Source: {filename} -->\n# {filename}\n\n{extracted_text}")

    if len(extracted_items) == 1:
        return jsonify({
            "success": True,
            "filename": extracted_items[0]["filename"],
            "content": extracted_items[0]["content"],
            "items": extracted_items,
            "count": 1
        })
    else:
        return jsonify({
            "success": True,
            "filename": f"batch_upload_{len(extracted_items)}_files.md",
            "content": "\n\n---\n\n".join(combined_texts),
            "items": extracted_items,
            "count": len(extracted_items)
        })

@standards_bp.route('/<path:file_id>', methods=['DELETE'])
def delete_standard(file_id):
    try:
        if file_id.isdigit():
            execute_write("DELETE FROM architecture_standards WHERE id=%s", (file_id,))
        else:
            filename = os.path.basename(file_id)
            title = filename.replace('.md', '').replace('_', ' ').replace('-', ' ').title()
            execute_write("DELETE FROM architecture_standards WHERE title=%s OR title=%s", (filename, title))
            file_path = os.path.join(KB_DIR, file_id)
            if os.path.exists(file_path):
                os.remove(file_path)
        return jsonify({"success": True, "message": "Standard deleted successfully"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500
